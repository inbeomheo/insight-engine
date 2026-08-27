"""
문서 텍스트 추출 서비스

PDF, DOCX, PPTX 파일에서 텍스트를 추출하여 콘텐츠 생성 파이프라인에 전달.
"""
import logging
import multiprocessing
import os
import tempfile
import time
import zipfile
from typing import Any, Dict

logger = logging.getLogger(__name__)

# 허용 MIME 타입
ALLOWED_MIME_TYPES = {
    'application/pdf': 'pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
}

# 파일 확장자 → MIME 매핑 (MIME 감지 실패 시 폴백)
EXTENSION_TO_MIME = {
    '.pdf': 'application/pdf',
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
}

MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
MAX_OOXML_ENTRIES = 2_048
MAX_OOXML_UNCOMPRESSED_BYTES = 64 * 1024 * 1024
MAX_OOXML_ENTRY_BYTES = 32 * 1024 * 1024
MAX_OOXML_COMPRESSION_RATIO = 200
MAX_PDF_PAGES = 200
MAX_EXTRACTED_TEXT_CHARS = 200_000
MAX_PDF_PROCESS_SECONDS = 20
MAX_PDF_WORKER_MEMORY_BYTES = 256 * 1024 * 1024


def _validate_ooxml_archive(file_path: str) -> None:
    """Reject malformed or excessively expanded DOCX/PPTX archives before parsing."""
    try:
        with zipfile.ZipFile(file_path) as archive:
            entries = archive.infolist()
            if len(entries) > MAX_OOXML_ENTRIES:
                raise ValueError('문서 내부 파일 수가 허용 한도를 초과합니다.')

            total_size = 0
            for entry in entries:
                if entry.is_dir():
                    continue
                normalized_name = entry.filename.replace('\\', '/')
                if normalized_name.startswith('/') or '..' in normalized_name.split('/'):
                    raise ValueError('문서 내부 경로가 안전하지 않습니다.')
                if entry.flag_bits & 0x1:
                    raise ValueError('암호화된 DOCX/PPTX 파일은 지원하지 않습니다.')
                if entry.file_size > MAX_OOXML_ENTRY_BYTES:
                    raise ValueError('문서 내부 파일 크기가 허용 한도를 초과합니다.')

                total_size += entry.file_size
                if total_size > MAX_OOXML_UNCOMPRESSED_BYTES:
                    raise ValueError('문서 압축 해제 크기가 허용 한도를 초과합니다.')

                if entry.file_size:
                    if entry.compress_size <= 0:
                        raise ValueError('비정상적인 문서 압축 구조입니다.')
                    ratio = entry.file_size / entry.compress_size
                    if ratio > MAX_OOXML_COMPRESSION_RATIO:
                        raise ValueError('문서 압축률이 허용 한도를 초과합니다.')
    except zipfile.BadZipFile as exc:
        raise ValueError('유효한 DOCX/PPTX 압축 문서가 아닙니다.') from exc


def extract_text(file_path: str, mime_type: str) -> Dict:
    """파일에서 텍스트를 추출합니다.

    Args:
        file_path: 파일 경로
        mime_type: MIME 타입 ('application/pdf' 또는 'application/vnd.openxmlformats-...')

    Returns:
        {"title": str, "content": str, "source_type": "document", "page_count": int}

    Raises:
        ValueError: 지원하지 않는 파일 형식이거나 텍스트 추출 실패
    """
    file_type = ALLOWED_MIME_TYPES.get(mime_type)
    if not file_type:
        raise ValueError(f'지원하지 않는 파일 형식입니다: {mime_type}')

    if file_type == 'pdf':
        return _extract_pdf(file_path)
    elif file_type == 'pptx':
        return _extract_pptx(file_path)
    else:
        return _extract_docx(file_path)


def _extract_pdf(file_path: str) -> Dict:
    """Extract an untrusted PDF in a time- and memory-bounded subprocess."""
    try:
        file_size = os.path.getsize(file_path)
    except OSError as exc:
        raise ValueError('PDF 파일을 읽을 수 없습니다.') from exc
    if file_size <= 0:
        raise ValueError('빈 PDF 파일입니다.')
    if file_size > MAX_FILE_SIZE:
        raise ValueError(f'파일 크기가 {MAX_FILE_SIZE // (1024 * 1024)}MB를 초과합니다.')

    context = multiprocessing.get_context('spawn')
    receiver, sender = context.Pipe(duplex=False)
    process = context.Process(
        target=_pdf_worker,
        args=(file_path, sender),
        name='pdf-extraction-worker',
        daemon=True,
    )
    started = False
    try:
        process.start()
        started = True
        sender.close()
        if not receiver.poll(MAX_PDF_PROCESS_SECONDS):
            raise ValueError('PDF 처리 시간이 허용 한도를 초과합니다.')
        try:
            status, payload = receiver.recv()
        except EOFError as exc:
            raise ValueError('PDF 안전 처리 프로세스가 비정상 종료되었습니다.') from exc
        if status == 'ok' and isinstance(payload, dict):
            return payload
        if status == 'error' and isinstance(payload, str):
            raise ValueError(payload)
        raise ValueError('PDF 안전 처리 프로세스가 올바른 결과를 반환하지 않았습니다.')
    finally:
        receiver.close()
        sender.close()
        if started:
            process.join(timeout=1)
            if process.is_alive():
                process.terminate()
                process.join(timeout=1)
            if process.is_alive() and hasattr(process, 'kill'):
                process.kill()
                process.join(timeout=1)


def _apply_pdf_worker_limits() -> None:
    """Apply OS limits when available; the parent timeout remains mandatory."""
    try:
        import resource
    except ImportError:
        # Windows and some macOS/container combinations do not expose every
        # rlimit. The isolated process and parent wall-clock timeout still apply.
        return
    limits = [
        (
            resource.RLIMIT_AS,
            (MAX_PDF_WORKER_MEMORY_BYTES, MAX_PDF_WORKER_MEMORY_BYTES),
        ),
        (
            resource.RLIMIT_CPU,
            (
                max(1, MAX_PDF_PROCESS_SECONDS - 2),
                max(1, MAX_PDF_PROCESS_SECONDS - 2),
            ),
        ),
    ]
    if hasattr(resource, 'RLIMIT_NOFILE'):
        limits.append((resource.RLIMIT_NOFILE, (64, 64)))
    for resource_kind, resource_limit in limits:
        try:
            resource.setrlimit(resource_kind, resource_limit)
        except (OSError, ValueError):
            continue


def _pdf_worker(file_path: str, sender) -> None:
    try:
        _apply_pdf_worker_limits()
        sender.send(('ok', _extract_pdf_payload(file_path)))
    except ValueError as exc:
        sender.send(('error', str(exc)))
    except BaseException:
        # Never serialize parser internals or paths back to the HTTP process.
        try:
            sender.send(('error', 'PDF를 안전하게 처리할 수 없습니다.'))
        except Exception:
            pass
    finally:
        sender.close()


def _extract_pdf_payload(file_path: str) -> Dict:
    """PDF parser body executed only inside the constrained worker."""
    import pypdf

    started_at = time.monotonic()
    try:
        with open(file_path, 'rb') as f:
            reader = pypdf.PdfReader(f)

            if reader.is_encrypted:
                raise ValueError('암호화된 PDF 파일은 지원하지 않습니다.')

            page_count = len(reader.pages)
            if page_count == 0:
                raise ValueError('빈 PDF 파일입니다.')
            if page_count > MAX_PDF_PAGES:
                raise ValueError(f'PDF 페이지 수는 최대 {MAX_PDF_PAGES}페이지입니다.')

            pages = []
            total_chars = 0
            for page in reader.pages:
                if time.monotonic() - started_at > MAX_PDF_PROCESS_SECONDS - 2:
                    raise ValueError('PDF 처리 시간이 허용 한도를 초과합니다.')
                text = page.extract_text()
                if text:
                    normalized = text.strip()
                    total_chars += len(normalized)
                    if total_chars > MAX_EXTRACTED_TEXT_CHARS:
                        raise ValueError('PDF 추출 텍스트가 허용 한도를 초과합니다.')
                    pages.append(normalized)

            content = '\n\n'.join(pages)
            if not content.strip():
                raise ValueError('PDF에서 텍스트를 추출할 수 없습니다. 이미지 기반 PDF일 수 있습니다.')

            # 제목: 메타데이터 또는 파일명
            title = ''
            metadata = reader.metadata
            if metadata and metadata.title:
                title = str(metadata.title)[:1_000]

            if not title:
                title = os.path.splitext(os.path.basename(file_path))[0]

            return {
                'title': title,
                'content': content,
                'source_type': 'document',
                'page_count': page_count,
            }
    except pypdf.errors.PdfReadError as e:
        raise ValueError(f'PDF 파일을 읽을 수 없습니다: {e}')


def _extract_docx(file_path: str) -> Dict:
    """DOCX에서 텍스트를 추출합니다."""
    import docx

    _validate_ooxml_archive(file_path)

    try:
        doc = docx.Document(file_path)
    except Exception as e:
        raise ValueError(f'DOCX 파일을 읽을 수 없습니다: {e}')

    parts = []

    # 본문 paragraph
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # 테이블 텍스트
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(' | '.join(cells))

    content = '\n\n'.join(parts)
    if not content.strip():
        raise ValueError('DOCX 파일에서 텍스트를 추출할 수 없습니다.')

    # 제목: core_properties 또는 파일명
    title = ''
    try:
        if doc.core_properties.title:
            title = doc.core_properties.title
    except Exception:
        pass

    if not title:
        title = os.path.splitext(os.path.basename(file_path))[0]

    return {
        'title': title,
        'content': content,
        'source_type': 'document',
        'page_count': len(doc.paragraphs),  # DOCX는 paragraph 수를 참고값으로
    }


def _extract_pptx(file_path: str) -> Dict:
    """PPTX에서 슬라이드별 제목, 본문, 발표자 노트를 추출합니다."""
    from pptx import Presentation

    _validate_ooxml_archive(file_path)

    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise ValueError(f'PPTX 파일을 읽을 수 없습니다: {e}')

    slides_text = [
        text for idx, slide in enumerate(prs.slides, 1)
        if (text := _extract_slide_text(slide, idx))
    ]

    content = '\n\n'.join(slides_text)
    if not content.strip():
        raise ValueError('PPTX 파일에서 텍스트를 추출할 수 없습니다.')

    # 제목: 첫 번째 슬라이드 제목 또는 파일명
    title = ''
    if prs.slides and prs.slides[0].shapes.title:
        title = prs.slides[0].shapes.title.text.strip()
    if not title:
        title = os.path.splitext(os.path.basename(file_path))[0]

    return {
        'title': title,
        'content': content,
        'source_type': 'document',
        'page_count': len(prs.slides),
    }


def _extract_slide_text(slide, idx: int) -> str:
    """단일 슬라이드에서 제목, 본문, 노트를 추출합니다."""
    parts = []

    if slide.shapes.title and slide.shapes.title.text.strip():
        parts.append(f'### {slide.shapes.title.text.strip()}')

    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)

    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            parts.append(f'> 노트: {notes}')

    if parts:
        return f'## 슬라이드 {idx}\n\n' + '\n\n'.join(parts)
    return ''


def extract_from_upload(file_storage: Any) -> Dict:
    """Flask FileStorage 객체에서 텍스트를 추출합니다.

    Args:
        file_storage: Flask request.files['file'] 객체

    Returns:
        extract_text()와 동일한 Dict

    Raises:
        ValueError: 파일 크기 초과, 지원하지 않는 형식 등
    """
    filename = file_storage.filename or ''
    ext = os.path.splitext(filename)[1].lower()

    # 라우트의 확장자 기반 검증과 동일하게 파일명 확장자를 우선합니다.
    mime_type = EXTENSION_TO_MIME.get(ext)
    if mime_type not in ALLOWED_MIME_TYPES:
        mime_type = file_storage.content_type

    if not mime_type or mime_type not in ALLOWED_MIME_TYPES:
        raise ValueError(f'지원하지 않는 파일 형식입니다. PDF, DOCX, PPTX 파일만 업로드 가능합니다.')

    # 임시 파일로 저장 후 추출
    suffix = ext or ('.pdf' if 'pdf' in mime_type else '.docx')
    tmp_fd, tmp_path = tempfile.mkstemp(suffix=suffix)
    try:
        file_storage.save(tmp_path)

        # 파일 크기 검증
        file_size = os.path.getsize(tmp_path)
        if file_size > MAX_FILE_SIZE:
            raise ValueError(f'파일 크기가 {MAX_FILE_SIZE // (1024 * 1024)}MB를 초과합니다.')
        if file_size == 0:
            raise ValueError('빈 파일입니다.')

        result = extract_text(tmp_path, mime_type)
        # 파일명에서 제목 보완
        if not result.get('title') or result['title'] == os.path.splitext(os.path.basename(tmp_path))[0]:
            result['title'] = os.path.splitext(filename)[0]
        return result
    finally:
        os.close(tmp_fd)
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)
