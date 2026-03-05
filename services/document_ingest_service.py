"""
문서 텍스트 추출 서비스

PDF, DOCX, PPTX 파일에서 텍스트를 추출하여 콘텐츠 생성 파이프라인에 전달.
"""
import logging
import os
import tempfile
from typing import Dict

logger = logging.getLogger(__name__)

# 허용 MIME 타입
ALLOWED_MIME_TYPES = {
    'application/pdf': 'pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
}

# 파일 확장자 → MIME 매핑 (MIME 감지 실패 시 폴백)
EXTENSION_TO_MIME = {
    '.pdf': 'application/pdf',
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
}

MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB


def extract_text(file_path: str, mime_type: str) -> Dict:
    """파일에서 텍스트를 추출합니다.

    Args:
        file_path: 파일 경로
        mime_type: MIME 타입 ('application/pdf' 또는 'application/vnd.openxmlformats-...')

    Returns:
        {"title": str, "content": str, "source_type": "document", "page_count": int}

    Raises:
        ValueError: 지원하지 않는 파일 형식이거나 텍스트 추출 실패
    """
    file_type = ALLOWED_MIME_TYPES.get(mime_type)
    if not file_type:
        raise ValueError(f'지원하지 않는 파일 형식입니다: {mime_type}')

    if file_type == 'pdf':
        return _extract_pdf(file_path)
    elif file_type == 'pptx':
        return _extract_pptx(file_path)
    else:
        return _extract_docx(file_path)


def _extract_pdf(file_path: str) -> Dict:
    """PDF에서 텍스트를 추출합니다."""
    import PyPDF2

    try:
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)

            if reader.is_encrypted:
                raise ValueError('암호화된 PDF 파일은 지원하지 않습니다.')

            page_count = len(reader.pages)
            if page_count == 0:
                raise ValueError('빈 PDF 파일입니다.')

            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text.strip())

            content = '\n\n'.join(pages)
            if not content.strip():
                raise ValueError('PDF에서 텍스트를 추출할 수 없습니다. 이미지 기반 PDF일 수 있습니다.')

            # 제목: 메타데이터 또는 파일명
            title = ''
            metadata = reader.metadata
            if metadata and metadata.title:
                title = metadata.title

            if not title:
                title = os.path.splitext(os.path.basename(file_path))[0]

            return {
                'title': title,
                'content': content,
                'source_type': 'document',
                'page_count': page_count,
            }
    except PyPDF2.errors.PdfReadError as e:
        raise ValueError(f'PDF 파일을 읽을 수 없습니다: {e}')


def _extract_docx(file_path: str) -> Dict:
    """DOCX에서 텍스트를 추출합니다."""
    import docx

    try:
        doc = docx.Document(file_path)
    except Exception as e:
        raise ValueError(f'DOCX 파일을 읽을 수 없습니다: {e}')

    parts = []

    # 본문 paragraph
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # 테이블 텍스트
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(' | '.join(cells))

    content = '\n\n'.join(parts)
    if not content.strip():
        raise ValueError('DOCX 파일에서 텍스트를 추출할 수 없습니다.')

    # 제목: core_properties 또는 파일명
    title = ''
    try:
        if doc.core_properties.title:
            title = doc.core_properties.title
    except Exception:
        pass

    if not title:
        title = os.path.splitext(os.path.basename(file_path))[0]

    return {
        'title': title,
        'content': content,
        'source_type': 'document',
        'page_count': len(doc.paragraphs),  # DOCX는 paragraph 수를 참고값으로
    }


def _extract_pptx(file_path: str) -> Dict:
    """PPTX에서 슬라이드별 제목, 본문, 발표자 노트를 추출합니다."""
    from pptx import Presentation

    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise ValueError(f'PPTX 파일을 읽을 수 없습니다: {e}')

    slides_text = []
    for idx, slide in enumerate(prs.slides, 1):
        parts = []

        # 슬라이드 제목
        if slide.shapes.title and slide.shapes.title.text.strip():
            parts.append(f'### {slide.shapes.title.text.strip()}')

        # 본문 텍스트 (제목 제외)
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)

        # 발표자 노트
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f'> 노트: {notes}')

        if parts:
            slides_text.append(f'## 슬라이드 {idx}\n\n' + '\n\n'.join(parts))

    content = '\n\n'.join(slides_text)
    if not content.strip():
        raise ValueError('PPTX 파일에서 텍스트를 추출할 수 없습니다.')

    # 제목: 첫 번째 슬라이드 제목 또는 파일명
    title = ''
    if prs.slides and prs.slides[0].shapes.title:
        title = prs.slides[0].shapes.title.text.strip()
    if not title:
        title = os.path.splitext(os.path.basename(file_path))[0]

    return {
        'title': title,
        'content': content,
        'source_type': 'document',
        'page_count': len(prs.slides),
    }


def extract_from_upload(file_storage) -> Dict:
    """Flask FileStorage 객체에서 텍스트를 추출합니다.

    Args:
        file_storage: Flask request.files['file'] 객체

    Returns:
        extract_text()와 동일한 Dict

    Raises:
        ValueError: 파일 크기 초과, 지원하지 않는 형식 등
    """
    filename = file_storage.filename or ''
    ext = os.path.splitext(filename)[1].lower()

    # MIME 타입 결정 (Content-Type 우선, 확장자 폴백)
    mime_type = file_storage.content_type
    if mime_type not in ALLOWED_MIME_TYPES:
        mime_type = EXTENSION_TO_MIME.get(ext)

    if not mime_type or mime_type not in ALLOWED_MIME_TYPES:
        raise ValueError(f'지원하지 않는 파일 형식입니다. PDF, DOCX, PPTX 파일만 업로드 가능합니다.')

    # 임시 파일로 저장 후 추출
    suffix = ext or ('.pdf' if 'pdf' in mime_type else '.docx')
    tmp_fd, tmp_path = tempfile.mkstemp(suffix=suffix)
    try:
        file_storage.save(tmp_path)

        # 파일 크기 검증
        file_size = os.path.getsize(tmp_path)
        if file_size > MAX_FILE_SIZE:
            raise ValueError(f'파일 크기가 {MAX_FILE_SIZE // (1024 * 1024)}MB를 초과합니다.')
        if file_size == 0:
            raise ValueError('빈 파일입니다.')

        result = extract_text(tmp_path, mime_type)
        # 파일명에서 제목 보완
        if not result.get('title') or result['title'] == os.path.splitext(os.path.basename(tmp_path))[0]:
            result['title'] = os.path.splitext(filename)[0]
        return result
    finally:
        os.close(tmp_fd)
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)
